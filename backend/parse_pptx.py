from flask import Flask, request, jsonify
from pptx import Presentation
import os
import logging
import re
from collections import OrderedDict
import spacy
from transformers import pipeline
from keybert import KeyBERT

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize Flask app
app = Flask(__name__)

# Load NLP models
try:
    nlp = spacy.load("en_core_web_sm")
except OSError:
    logger.error("spaCy model 'en_core_web_sm' not found. Installing now...")
    spacy.cli.download("en_core_web_sm")
    nlp = spacy.load("en_core_web_sm")

# Initialize summarizer and keyword extractor
try:
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    kw_model = KeyBERT()
except Exception as e:
    logger.error(f"Failed to initialize NLP models: {str(e)}")
    raise

# Utility functions
def clean_text(text: str) -> str:
    """Clean and normalize text for readability, handling bullet points and fragments."""
    if not text:
        return ""
    cleaned = re.sub(r'^\s*-?\s*', '', text, flags=re.MULTILINE)
    cleaned = re.sub(r'https?://\S+', '', cleaned)
    cleaned = re.sub(r'\b\d{1,2}[^\w\s]\d{1,2}[^\w\s]\d{2,4}\b', '', cleaned)
    cleaned = re.sub(r'\b\d+\b', '', cleaned)
    cleaned = re.sub(r'\s*\(ex\.\s*[^\)]*\)', '', cleaned)
    cleaned = re.sub(r'\s*\(money,\s*p\d+\)', '', cleaned)
    cleaned = re.sub(r'\b(quiz|interactive|task|upload|canvas|source|video|slide|page)\b', '', cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    return cleaned

def is_complete_sentence(sentence: str, nlp_model) -> bool:
    """Check if a sentence is grammatically complete using spaCy."""
    doc = nlp_model(sentence)
    if not doc or not doc.sents:
        return False
    has_subject = any(token.dep_ in ['nsubj', 'nsubjpass'] for token in doc)
    has_verb = any(token.pos_ == 'VERB' and token.dep_ in ['ROOT', 'aux'] for token in doc)
    return has_subject and has_verb and len(doc) >= 5

def extract_notes(text: str, images: list = None) -> dict:
    """Extract readable notes with only complete sentences for guiding questions and notes."""
    notes = {
        "guiding_questions": [],
        "notes": [],
        "image_references": images or []
    }

    cleaned_text = clean_text(text)
    if not cleaned_text and not images:
        return notes

    doc = nlp(cleaned_text)
    for sentence in doc.sents:
        sentence_text = sentence.text.strip()
        if sentence_text and sentence_text.endswith('?'):
            if is_complete_sentence(sentence_text, nlp) and len(sentence_text.split()) >= 5 and sentence_text not in notes["guiding_questions"]:
                notes["guiding_questions"].append(sentence_text.capitalize())
        elif sentence_text:
            if is_complete_sentence(sentence_text, nlp) and len(sentence_text.split()) >= 5 and sentence_text not in notes["notes"]:
                notes["notes"].append(sentence_text.capitalize())

    return notes

# API Routes
@app.route('/parse-pptx', methods=['POST'])
def parse_pptx():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if not file or file.filename == '':
        return jsonify({'error': 'No valid file provided'}), 400

    file_path = os.path.join('uploads', file.filename)
    os.makedirs('uploads', exist_ok=True)

    try:
        logger.info(f"Received file: {file.filename}")
        file.save(file_path)
        prs = Presentation(file_path)
        all_notes = {
            'guiding_questions': [],
            'definitions': [],
            'specific_topics': {},
            'key_phrases': [],
            'summary': '',
            'notes': [],
            'image_references': []
        }

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = {'title': '', 'notes': {'raw_text': ''}, 'images': []}
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    if not slide_content['title']:
                        slide_content['title'] = shape.text.strip()
                    else:
                        slide_content['notes']['raw_text'] += shape.text.strip() + ' '
                if shape.shape_type == 13:  # Picture type
                    slide_content['images'].append(f"Image_{slide_num}_{shape.shape_id}")

            if slide_content['notes']['raw_text'] or slide_content['images']:
                notes = extract_notes(slide_content['notes']['raw_text'], slide_content['images'])
                all_notes['guiding_questions'].extend(notes["guiding_questions"])
                all_notes['notes'].extend(notes["notes"])
                all_notes['image_references'].extend(notes["image_references"])

        final_notes = {
            'guiding_questions': list(OrderedDict.fromkeys(all_notes['guiding_questions'])),
            'definitions': [],
            'specific_topics': {},
            'key_phrases': [],
            'summary': '',
            'notes': list(OrderedDict.fromkeys(all_notes['notes'])),
            'image_references': sorted(list(set(all_notes['image_references'])))
        }

        logger.info(f"Returning notes: {final_notes}")
        if not final_notes['guiding_questions'] and not final_notes['notes']:
            logger.warning("No notes extracted—check filtering or input data.")
        os.remove(file_path)
        return jsonify({'topics': {'All Notes': final_notes}})

    except Exception as e:
        logger.error(f"Error parsing PowerPoint file: {str(e)}")
        return jsonify({'error': f"Failed to parse file: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=6000)